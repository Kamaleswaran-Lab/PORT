"""
make_adjudication.py
-----------------------
Step 2 (Audit-refined) 최종 결과 기반 adjudication 자료 생성.

Outputs:
  iod2_adjudication_50.csv       — 의사 리뷰용 50개 노트 (필터 유형별 층화)
  iod2_adjudication_process.pptx — 전체 과정 요약 프레젠테이션 (Calibri, 흰 배경)

Usage:
    conda activate ethos
    python datapreprocessing/meds_scripts/make_adjudication.py
"""

import re
import sys
import csv
import random
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = Path("/path/to/CHD_MEDS/outcome")

# ── copy all filter patterns from iod_to_outcome.py ─────────────────────────
exec(open("datapreprocessing/meds_scripts/iod_to_outcome.py").read().split("def read_rpt")[0])


# ── classify each note with reason label ────────────────────────────────────
def classify_note(text: str):
    """Returns (kept: bool, reason: str)."""
    t = text
    keep = True
    reason = "passed"

    epi_excluded = False
    help_excl = False
    icd_excl = False

    epi_rows = bool(re.search(r"\bepi", t, re.IGNORECASE))
    if epi_rows and re.search(_IOD2_EPI_EXCLUDE, t):
        keep = False; reason = "epi_FP"; epi_excluded = True

    echo_rows = bool(re.search(r"\becho\b", t, re.IGNORECASE))
    if echo_rows and re.search(_IOD2_ECHO_ROUTINE, t):
        keep = False; reason = "echo_routine" if reason == "passed" else reason
    if echo_rows and re.search(_IOD2_ECHO_LOGISTICAL, t):
        keep = False; reason = "echo_logistical" if reason == "passed" else reason

    help_rows = bool(re.search(r"\bhelp\b", t, re.IGNORECASE))
    if help_rows and not re.search(_IOD2_HELP_EMERGENCY, t):
        keep = False; reason = "help_not_emergency" if reason == "passed" else reason
        help_excl = True

    brady_rows = bool(re.search(r"\bbrady\b", t, re.IGNORECASE))
    is_dr_brady = brady_rows and bool(re.search(_IOD2_BRADY_EXCLUDE, t))
    has_clinical_brady = bool(re.search(r"bradycardia|bradycardic", t, re.IGNORECASE))
    if is_dr_brady and not has_clinical_brady:
        keep = False; reason = "dr_brady_name" if reason == "passed" else reason

    has_other_iod_signal = bool(re.search(
        r"epi(?:nephrine)?\s+(?:given|IV|drip|bolus|gtt|\d+\s*mcg)"
        r"|vasopressor|compressions|CPR|pulseless|code\s+(?:blue|cart)"
        r"|laryngospasm|bronchospasm|desaturat(?:ion|ed|ing)"
        r"|arrest\b|emergent|\bbradycardia\b|unable\s+to\s+ventilate"
        r"|atropine|glycopyrrolate|succinylcholine|\bsux\b",
        t, re.IGNORECASE
    ))
    brady_negated = bool(re.search(_IOD2_BRADY_NEGATION, t))
    if brady_negated and not has_other_iod_signal:
        keep = False; reason = "brady_negation" if reason == "passed" else reason

    defib_rows = bool(re.search(r"\bdefib", t, re.IGNORECASE))
    if defib_rows and re.search(_IOD2_DEFIB_EXCLUDE, t):
        keep = False; reason = "defib_test" if reason == "passed" else reason

    hypo_rows = bool(re.search(r"hypotension", t, re.IGNORECASE))
    if hypo_rows and re.search(_IOD2_HYPOTENSION_EXCLUDE, t):
        keep = False; reason = "planned_hypotension" if reason == "passed" else reason

    phenyl_rows = bool(re.search(r"phenylephrine", t, re.IGNORECASE))
    if phenyl_rows and re.search(_IOD2_PHENYL_OPHTHALMIC_EXCLUDE, t):
        keep = False; reason = "ophthalmic_phenylephrine" if reason == "passed" else reason

    nasal_epi = bool(re.search(r"\bepi", t, re.IGNORECASE))
    if nasal_epi and re.search(_IOD2_NASAL_EPI_EXCLUDE, t):
        keep = False; reason = "intranasal_topical_epi" if reason == "passed" else reason

    icd_rows = bool(re.search(
        r"\bICD\b|defibrillat|induced?\s+(?:VT|V-?fib)|intentionally\s+fibrillat", t, re.IGNORECASE
    ))
    if icd_rows and re.search(_IOD2_ICD_TEST_EXCLUDE, t):
        keep = False; reason = "icd_device_test" if reason == "passed" else reason
        icd_excl = True

    arrest_rows = bool(re.search(r"\barrest\b", t, re.IGNORECASE))
    if arrest_rows and re.search(_IOD2_CIRC_ARREST_EXCLUDE, t):
        keep = False; reason = "circ_arrest_bypass" if reason == "passed" else reason

    code_rows = bool(re.search(r"\bcode\b", t, re.IGNORECASE))
    if code_rows and re.search(_IOD2_CODE_FACILITY_EXCLUDE, t):
        keep = False; reason = "code_facility" if reason == "passed" else reason

    # Overrides
    if re.search(_IOD2_ALWAYS_KEEP, t):
        keep = True; reason = "override_always_keep"
    if epi_excluded and re.search(_IOD2_EPI_OVERRIDE, t):
        keep = True; reason = "override_epi_systemic"
    if icd_excl and re.search(_IOD2_ICD_FAIL_OVERRIDE, t):
        keep = True; reason = "override_icd_fail"
    if help_excl and re.search(_IOD2_HELP_CLINICAL_OVERRIDE, t):
        keep = True; reason = "override_help_clinical"

    return keep, reason


# ── friendly label names for PPTX / CSV ─────────────────────────────────────
REASON_LABELS = {
    "epi_FP":                   "Local/topical epinephrine (not IV)",
    "echo_routine":             "Routine echo start/stop marker",
    "echo_logistical":          "Echo scheduling/logistics note",
    "help_not_emergency":       "Routine 'help' (non-emergency)",
    "dr_brady_name":            "Dr. Brady (surgeon's name)",
    "brady_negation":           "'No bradycardia' (negation)",
    "defib_test":               "Defibrillator equipment test",
    "planned_hypotension":      "Planned/deliberate hypotension",
    "ophthalmic_phenylephrine": "Ophthalmic phenylephrine (eye drops)",
    "intranasal_topical_epi":   "Intranasal/topical epinephrine",
    "icd_device_test":          "ICD device functional test",
    "circ_arrest_bypass":       "Circulatory arrest during CPB (DHCA)",
    "code_facility":            "Facility code (fire/security, not cardiac)",
    "override_always_keep":     "KEPT — unambiguous emergency",
    "override_epi_systemic":    "KEPT — systemic IV epinephrine confirmed",
    "override_icd_fail":        "KEPT — ICD conversion failed (emergency)",
    "override_help_clinical":   "KEPT — hemodynamic collapse confirmed",
    "passed":                   "KEPT — passed all filters",
}


def build_adjudication_csv(an_events):
    """Build stratified 50-note adjudication CSV from removed notes."""
    qn = an_events[an_events["Event"].fillna("").str.strip() == "Quick Note"].copy()
    txt = qn["Event Comment"].fillna("")
    matched = qn[txt.str.contains(_IOD2_PATTERN, regex=True)].copy()

    records = []
    for _, row in matched.iterrows():
        text = str(row["Event Comment"]) if pd.notna(row["Event Comment"]) else ""
        kept, reason = classify_note(text)
        records.append({
            "C MRN": row.get("C MRN", ""),
            "Recorded Time": row.get("Recorded Time", ""),
            "Quick Note Text": text,
            "kept": kept,
            "reason": reason,
            "filter_label": REASON_LABELS.get(reason, reason),
        })

    df = pd.DataFrame(records)
    removed = df[~df["kept"]].copy()

    # Stratified sample: proportional to count, min 2 per category, total ~50
    target = 50
    reason_counts = removed["reason"].value_counts()

    # Priority reasons (most informative for adjudication)
    # Skip very high-n epi_FP (already well-characterised); focus on borderline categories
    priority_order = [
        "circ_arrest_bypass",       # new Step 2 — clinical significance question
        "echo_logistical",          # new Step 2 — was this truly non-IoD?
        "planned_hypotension",      # borderline — was it really planned?
        "help_not_emergency",       # high FN risk
        "echo_routine",             # established
        "dr_brady_name",            # name vs clinical
        "ophthalmic_phenylephrine", # topical vs systemic
        "icd_device_test",          # test vs emergency
        "defib_test",               # test vs emergency
        "brady_negation",           # new
        "code_facility",            # new
        "intranasal_topical_epi",   # established
        "epi_FP",                   # established but large — add a few
    ]

    samples = []
    remaining = target
    for reason in priority_order:
        if reason not in reason_counts.index:
            continue
        n_avail = reason_counts[reason]
        n_take = min(n_avail, max(2, min(6, remaining // max(1, len(priority_order) - priority_order.index(reason)))))
        if n_take > 0:
            sub = removed[removed["reason"] == reason].sample(min(n_take, n_avail), random_state=42)
            samples.append(sub)
            remaining -= len(sub)
        if remaining <= 0:
            break

    adjud = pd.concat(samples).reset_index(drop=True)
    adjud["row_num"] = range(1, len(adjud) + 1)
    adjud["Is_IoD_clinical"] = ""  # physician fills: Yes / No / Unclear
    adjud["Notes"] = ""            # physician comments

    out_cols = ["row_num", "C MRN", "Recorded Time", "Quick Note Text",
                "filter_label", "Is_IoD_clinical", "Notes"]
    adjud[out_cols].to_csv(OUTPUT_DIR / "iod2_adjudication_50.csv", index=False, encoding="utf-8-sig")
    print(f"Saved adjudication CSV: {len(adjud)} rows")
    print(adjud["reason"].value_counts().to_string())
    return df, adjud


# ── PPTX builder ─────────────────────────────────────────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x30, 0x30, 0x30)
ACCENT = RGBColor(0x1F, 0x5C, 0x99)  # dark blue — subtle

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs):
    """Add a completely blank slide (no placeholders)."""
    blank_layout = prs.slide_layouts[6]  # layout index 6 = blank
    return prs.slides.add_slide(blank_layout)


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, italic=False,
                color=BLACK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=14, bold=False, italic=False,
             color=BLACK, align=PP_ALIGN.LEFT, indent=0, font_name="Calibri"):
    p = tf.add_paragraph()
    p.alignment = align
    if indent:
        p.level = indent
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def set_slide_white(slide):
    """Set slide background to white."""
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def build_pptx(df_all, adjud_df):
    prs = new_prs()

    # ── Slide 1: Title ──────────────────────────────────────────────────────
    sl = blank_slide(prs)
    set_slide_white(sl)
    add_textbox(sl, "IoD2 Label Refinement",
                0.6, 2.2, 12, 1.0, font_size=36, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
    add_textbox(sl, "Physician Adjudication — Quick Note Filter Validation",
                0.6, 3.3, 12, 0.7, font_size=22, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT)
    add_textbox(sl, "CHOA Pediatric Cardiac Surgery · IoD Prediction Project · 2026",
                0.6, 4.2, 12, 0.5, font_size=14, italic=True, color=DARK_GRAY, align=PP_ALIGN.LEFT)

    # ── Slide 2: What is IoD2? ───────────────────────────────────────────────
    sl = blank_slide(prs)
    set_slide_white(sl)
    add_textbox(sl, "What is IoD2?", 0.6, 0.35, 12, 0.6, font_size=26, bold=True, color=ACCENT)
    txBox = sl.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.8))
    tf = txBox.text_frame; tf.word_wrap = True
    add_para(tf, "Intraoperative Deterioration component 2 (IoD2) is defined as:", 15, bold=True, color=DARK_GRAY)
    add_para(tf, "  → The presence of a Quick Note in the anesthesia record during the OR window", 14, color=BLACK)
    add_para(tf, "     that contains one or more keywords signaling hemodynamic or respiratory deterioration.", 14, color=BLACK)
    add_para(tf, "", 6)
    add_para(tf, "Keywords include:  CPR, arrest, chest compression, epinephrine, dopamine, norepinephrine,", 13, italic=True, color=DARK_GRAY)
    add_para(tf, "                   phenylephrine, hypotension, shock, defib, ECMO, SVT, V-tach, V-fib, brady,", 13, italic=True, color=DARK_GRAY)
    add_para(tf, "                   emergently, help, echo, cardiac output, arrhythmia, and others.", 13, italic=True, color=DARK_GRAY)
    add_para(tf, "", 6)
    add_para(tf, "Problem:", 15, bold=True, color=DARK_GRAY)
    add_para(tf, "  Many keyword matches are FALSE POSITIVES — the keyword appears in a routine context,", 14, color=BLACK)
    add_para(tf, "  not as a sign of deterioration.", 14, color=BLACK)
    add_para(tf, "", 5)
    add_para(tf, "  Example:  'epinephrine' in 'Surgeon injects 1% lido with epi 1:100K' = local anesthetic", 13, italic=True, color=DARK_GRAY)
    add_para(tf, "            'echo' in 'Echo start' = routine TEE monitoring marker", 13, italic=True, color=DARK_GRAY)
    add_para(tf, "            'arrest' in 'circ arrest' = planned deep hypothermic circulatory arrest (DHCA)", 13, italic=True, color=DARK_GRAY)

    # ── Slide 3: Refinement Process ─────────────────────────────────────────
    sl = blank_slide(prs)
    set_slide_white(sl)
    add_textbox(sl, "Two-Stage Iterative Label Refinement", 0.6, 0.35, 12, 0.6, font_size=26, bold=True, color=ACCENT)

    n_matched = len(df_all)
    n_kept_s1 = 3977   # Step 1 kept (historical)
    n_removed_s1 = n_matched - n_kept_s1
    n_kept_s2 = len(df_all[df_all["kept"]])
    n_removed_s2 = n_matched - n_kept_s2
    iod2_orig = 3502; iod2_s1 = 1785; iod2_s2 = 1651
    comp_orig = 3808; comp_s1 = 2106; comp_s2 = 1973

    txBox = sl.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.8))
    tf = txBox.text_frame; tf.word_wrap = True
    add_para(tf, f"Starting point:  {n_matched:,} Quick Notes matched broad IoD2 keyword pattern", 14, bold=True, color=DARK_GRAY)
    add_para(tf, "", 5)
    add_para(tf, "Step 1 — Rule-based keyword exclusion (completed 2026-03-29)", 15, bold=True, color=ACCENT)
    add_para(tf, f"  Filters:  epi/local-anesthetic · echo routine start/stop · help (non-emergency) · Dr. Brady name", 13, color=BLACK)
    add_para(tf, f"            defib equipment test · planned hypotension · ophthalmic phenylephrine · intranasal epi", 13, color=BLACK)
    add_para(tf, f"  Result:   {n_kept_s1:,} notes kept  ({n_removed_s1:,} removed,  {n_removed_s1/n_matched*100:.1f}%)", 13, bold=True, color=DARK_GRAY)
    add_para(tf, f"  IoD2: {iod2_orig:,} → {iod2_s1:,} encounters   Composite IoD: {comp_orig:,} → {comp_s1:,}   Positive rate: 1.8% → 1.0%", 13, color=BLACK)
    add_para(tf, "", 5)
    add_para(tf, "Step 2 — Full-dataset audit & refinement (completed 2026-03-29)", 15, bold=True, color=ACCENT)
    add_para(tf, f"  Method:   All {n_matched:,} matched notes classified; FP and FN patterns systematically identified", 13, color=BLACK)
    add_para(tf, f"  New FP filters added:  circulatory arrest during CPB (DHCA) · echo scheduling/logistics notes", 13, color=BLACK)
    add_para(tf, f"                         facility codes (code red) · 'no bradycardia' negation", 13, color=BLACK)
    add_para(tf, f"  New FN overrides added: failed ICD conversion · hemodynamic collapse with 'help' language · IV epi", 13, color=BLACK)
    add_para(tf, f"  Result:   {n_kept_s2:,} notes kept  ({n_removed_s2:,} removed,  {n_removed_s2/n_matched*100:.1f}%)", 13, bold=True, color=DARK_GRAY)
    add_para(tf, f"  IoD2: {iod2_s1:,} → {iod2_s2:,} encounters   Composite IoD: {comp_s1:,} → {comp_s2:,}   Positive rate: 1.0% → 0.94%", 13, color=BLACK)

    # ── Slide 4: Filter category breakdown ─────────────────────────────────
    sl = blank_slide(prs)
    set_slide_white(sl)
    add_textbox(sl, "Removed Notes — Filter Category Breakdown", 0.6, 0.35, 12, 0.6, font_size=26, bold=True, color=ACCENT)

    removed_df = df_all[~df_all["kept"]]
    rc = removed_df["reason"].value_counts()
    total_removed = len(removed_df)

    txBox = sl.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.8))
    tf = txBox.text_frame; tf.word_wrap = True
    add_para(tf, f"Total removed:  {total_removed:,} notes  ({total_removed/n_matched*100:.1f}% of matched)", 14, bold=True, color=DARK_GRAY)
    add_para(tf, "", 4)

    step2_new = {"circ_arrest_bypass", "echo_logistical", "code_facility", "brady_negation"}
    rows = [
        ("epi_FP",                   "Local/topical epi (surgeon inject, epidural, Epic, ratio)",  "Step 1"),
        ("help_not_emergency",       "Routine 'help' without emergency call language",             "Step 1"),
        ("ophthalmic_phenylephrine", "Phenylephrine eye drops/swabs by surgeon",                  "Step 1"),
        ("echo_routine",             "Echo standalone markers (start/stop/probe)",                 "Step 1"),
        ("dr_brady_name",            "Dr. Brady (surgeon's name, not bradycardia)",                "Step 1"),
        ("planned_hypotension",      "Deliberate/planned surgical hypotension",                    "Step 1"),
        ("icd_device_test",          "ICD device functional test (induced VF for testing)",        "Step 1"),
        ("defib_test",               "Defibrillator equipment test",                               "Step 1"),
        ("intranasal_topical_epi",   "Intranasal epi (nasal pledgets, topical vasoconstruction)",  "Step 1"),
        ("circ_arrest_bypass",       "Circulatory arrest during CPB — DHCA (planned, not IoD)",   "Step 2 "),
        ("echo_logistical",          "Echo scheduling/logistics (waiting for echo, echo team)",   "Step 2 "),
        ("code_facility",            "Facility code: code red/orange/silver (not cardiac)",        "Step 2 "),
        ("brady_negation",           "'No bradycardia' negation without other IoD signal",         "Step 2 "),
    ]
    for reason, label, step in rows:
        n = rc.get(reason, 0)
        pct = n / total_removed * 100 if total_removed else 0
        step_tag = f"  [{step}]"
        is_new = "" in step
        color = ACCENT if is_new else DARK_GRAY
        add_para(tf, f"  {label:<56}  n = {n:4d}  ({pct:4.1f}%)  {step_tag}", 11,
                 bold=is_new, color=color)

    # ── Slide 5: New filters explanation (DHCA + echo logistical) ──────────
    sl = blank_slide(prs)
    set_slide_white(sl)
    add_textbox(sl, "New in Step 2: Cardiac Surgery–Specific Patterns", 0.6, 0.35, 12, 0.6, font_size=26, bold=True, color=ACCENT)
    txBox = sl.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.8))
    tf = txBox.text_frame; tf.word_wrap = True

    add_para(tf, "1.  Circulatory Arrest During CPB (DHCA) — 29 notes removed", 15, bold=True, color=DARK_GRAY)
    add_para(tf, "  Cardiac surgery often uses Deep Hypothermic Circulatory Arrest (DHCA):", 13, color=BLACK)
    add_para(tf, "  the heart is deliberately stopped to allow bloodless repair of complex lesions.", 13, color=BLACK)
    add_para(tf, "  This is a PLANNED step, not a deterioration event.", 13, color=BLACK)
    add_para(tf, "", 4)
    add_para(tf, "  Examples removed:  'CPB pump off, Circ arrest'  ·  'circ arrest 12min'", 12, italic=True, color=DARK_GRAY)
    add_para(tf, "                     'Cardioplegia given during brief period of circ arrest'", 12, italic=True, color=DARK_GRAY)
    add_para(tf, "  Kept (override):   Any circ arrest note that ALSO documents CPR or pulseless state", 12, italic=True, color=BLACK)
    add_para(tf, "", 8)

    add_para(tf, "2.  Echo Scheduling / Logistics Notes — 22 notes removed", 15, bold=True, color=DARK_GRAY)
    add_para(tf, "  Step 1 already removed standalone 'Echo start / Echo stop / Echo probe' markers.", 13, color=BLACK)
    add_para(tf, "  Step 2 extended this to broader scheduling language without clinical findings:", 13, color=BLACK)
    add_para(tf, "", 4)
    add_para(tf, "  Examples removed:  'Waiting for Echo Tech'  ·  'Echo team present for transthoracic ECHO'", 12, italic=True, color=DARK_GRAY)
    add_para(tf, "                     'Doing echo in lab'  ·  'Begin EKG and ECHO, in conjunction with ABR'", 12, italic=True, color=DARK_GRAY)
    add_para(tf, "  Clinical echo notes kept:  'Chest reopened because of diminished RV function by echo'", 12, italic=True, color=BLACK)
    add_para(tf, "                             'Echo shows severe MR with LVOT obstruction'", 12, italic=True, color=BLACK)

    # ── Slide 6: FN recovery examples ───────────────────────────────────────
    sl = blank_slide(prs)
    set_slide_white(sl)
    add_textbox(sl, "New in Step 2: False Negative Recovery (Override Rules)", 0.6, 0.35, 12, 0.6, font_size=26, bold=True, color=ACCENT)
    txBox = sl.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.8))
    tf = txBox.text_frame; tf.word_wrap = True

    add_para(tf, "These notes were INCORRECTLY REMOVED by Step 1 filters — now restored:", 14, bold=True, color=DARK_GRAY)
    add_para(tf, "", 6)
    add_para(tf, "A.  ICD Failure Override  (1 note restored)", 14, bold=True, color=DARK_GRAY)
    add_para(tf, '  Note: "Internal defibrillation for V. Fibrillation, multiple defibrillations at 50J', 13, italic=True, color=BLACK)
    add_para(tf, '        without successful conversion of rhythm"', 13, italic=True, color=BLACK)
    add_para(tf, '  Why removed:  matched ICD device test pattern ("internal defibrillation")', 12, color=DARK_GRAY)
    add_para(tf, '  Why restored: "without successful conversion" signals failed rescue — true emergency', 12, bold=True, color=DARK_GRAY)
    add_para(tf, "", 8)
    add_para(tf, "B.  Help + Hemodynamic Collapse Override  (1 note restored)", 14, bold=True, color=DARK_GRAY)
    add_para(tf, '  Note: "Desaturation to 80% and falling. Presumed bronchospasm.', 13, italic=True, color=BLACK)
    add_para(tf, '        Call for additional help and administered first dose epinephrine IV"', 13, italic=True, color=BLACK)
    add_para(tf, '  Why removed:  "call for additional help" did not match old emergency-help patterns', 12, color=DARK_GRAY)
    add_para(tf, '  Why restored: explicit desaturation + IV epinephrine = true emergency', 12, bold=True, color=DARK_GRAY)
    add_para(tf, "", 8)
    add_para(tf, "C.  Help Emergency Pattern Expanded", 14, bold=True, color=DARK_GRAY)
    add_para(tf, '  Added: "call for additional help" / "called for additional help" to emergency patterns', 13, color=BLACK)

    # ── Slide 7: Final statistics ────────────────────────────────────────────
    sl = blank_slide(prs)
    set_slide_white(sl)
    add_textbox(sl, "Final Label Statistics (Step 2, 2026-03-29)", 0.6, 0.35, 12, 0.6, font_size=26, bold=True, color=ACCENT)
    txBox = sl.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.8))
    tf = txBox.text_frame; tf.word_wrap = True

    add_para(tf, "Quick Note Filter Pipeline Summary:", 15, bold=True, color=DARK_GRAY)
    add_para(tf, "", 4)
    rows2 = [
        ("Total Quick Notes in dataset",              "104,841"),
        ("Matched IoD2 keyword pattern (broad)",      f"{n_matched:,}"),
        ("Removed by FP filters (Step 1 + Step 2)",   f"{n_removed_s2:,}  ({n_removed_s2/n_matched*100:.1f}%)"),
        ("Final kept (IoD2-positive notes)",           f"{n_kept_s2:,}  ({n_kept_s2/n_matched*100:.1f}%)"),
    ]
    for label, val in rows2:
        add_para(tf, f"  {label:<52}  {val}", 13, color=BLACK)

    add_para(tf, "", 8)
    add_para(tf, "IoD Label Prevalence (210,274 total encounters):", 15, bold=True, color=DARK_GRAY)
    add_para(tf, "", 4)
    rows3 = [
        ("IoD1 — CPR during OR window",                "85",    "0.04%"),
        ("IoD2 — Quick Note keywords (Step 2 final)",  "1,651", "0.78%"),
        ("IoD3 — IV vasoactive bolus during OR",       "142",   "0.07%"),
        ("IoD4 — Vasoactive infusion started/escalated","152",  "0.07%"),
        ("IoD5 — Arterial line after incision",        "93",    "0.04%"),
        ("IoD COMPOSITE (IoD1–5)",                     "1,973", "0.94%"),
    ]
    for label, n_str, pct in rows3:
        bold = "COMPOSITE" in label
        col = ACCENT if bold else BLACK
        add_para(tf, f"  {label:<52}  n = {n_str:>6}   ({pct})", 13, bold=bold, color=col)

    add_para(tf, "", 6)
    add_para(tf, "Compared to original (pre-filter): IoD2: 3,502 → 1,651  (−53%)   Composite: 3,808 → 1,973  (−48%)", 13, italic=True, color=DARK_GRAY)

    # ── Slide 8: Adjudication instructions ──────────────────────────────────
    sl = blank_slide(prs)
    set_slide_white(sl)
    add_textbox(sl, "Physician Adjudication — Instructions", 0.6, 0.35, 12, 0.6, font_size=26, bold=True, color=ACCENT)
    txBox = sl.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.8))
    tf = txBox.text_frame; tf.word_wrap = True

    add_para(tf, "File provided:  iod2_adjudication_50.csv  (50 Quick Notes, all currently REMOVED by filters)", 14, bold=True, color=DARK_GRAY)
    add_para(tf, "", 5)
    add_para(tf, "Columns:", 14, bold=True, color=DARK_GRAY)
    add_para(tf, "  Quick Note Text     — the free-text note from the anesthesia record", 13, color=BLACK)
    add_para(tf, "  filter_label        — why the filter algorithm removed this note", 13, color=BLACK)
    add_para(tf, "  Is_IoD_clinical     — YOUR RESPONSE: Yes / No / Unclear", 13, bold=True, color=DARK_GRAY)
    add_para(tf, "  Notes               — optional: brief comment if needed", 13, color=BLACK)
    add_para(tf, "", 6)
    add_para(tf, "Definition of 'IoD' for this review:", 14, bold=True, color=DARK_GRAY)
    add_para(tf, "  YES  — the note describes an acute hemodynamic or respiratory deterioration during surgery", 13, color=BLACK)
    add_para(tf, "         that required emergent intervention (vasopressors, CPR, intubation, cardioversion, etc.)", 13, color=BLACK)
    add_para(tf, "  NO   — the note describes a routine event (local anesthetic injection, planned procedure,", 13, color=BLACK)
    add_para(tf, "         transport update, equipment note, name mention, etc.)", 13, color=BLACK)
    add_para(tf, "  UNCLEAR — insufficient information to judge, or borderline", 13, color=BLACK)
    add_para(tf, "", 6)
    add_para(tf, "Goal:  Estimate how many 'removed' notes are actually true IoD events (false negatives).", 13, italic=True, color=DARK_GRAY)
    add_para(tf, "       Results will guide final filter calibration before model retraining.", 13, italic=True, color=DARK_GRAY)

    out_path = OUTPUT_DIR / "iod2_adjudication_process.pptx"
    prs.save(str(out_path))
    print(f"Saved PPTX: {out_path}")


# ── main ─────────────────────────────────────────────────────────────────────
def main():
    FILE_AN_EVENTS = (
        "/path/to/CHOA_RAW_TABLES/"
        "CHOA_DATA_Tables_CHD/DR15201_AN_Events.rpt"
    )
    print("Loading AN_Events...")
    an_events = pd.read_csv(
        FILE_AN_EVENTS, delimiter="|", encoding="utf-8-sig",
        encoding_errors="replace", on_bad_lines="skip", low_memory=False
    )

    print("Classifying notes with Step 2 filters...")
    qn = an_events[an_events["Event"].fillna("").str.strip() == "Quick Note"].copy()
    txt = qn["Event Comment"].fillna("")
    matched = qn[txt.str.contains(_IOD2_PATTERN, regex=True)].copy()

    records = []
    for _, row in matched.iterrows():
        text = str(row["Event Comment"]) if pd.notna(row["Event Comment"]) else ""
        kept, reason = classify_note(text)
        records.append({
            "C MRN": row.get("C MRN", ""),
            "Recorded Time": row.get("Recorded Time", ""),
            "Quick Note Text": text,
            "kept": kept,
            "reason": reason,
            "filter_label": REASON_LABELS.get(reason, reason),
        })
    df_all = pd.DataFrame(records)

    n_kept = df_all["kept"].sum()
    n_rem = (~df_all["kept"]).sum()
    print(f"Matched: {len(df_all):,}  Kept: {n_kept:,}  Removed: {n_rem:,}")

    print("\nBuilding adjudication CSV...")
    _, adjud_df = build_adjudication_csv(an_events)

    print("\nBuilding PPTX...")
    build_pptx(df_all, adjud_df)
    print("\nDone.")


if __name__ == "__main__":
    main()
